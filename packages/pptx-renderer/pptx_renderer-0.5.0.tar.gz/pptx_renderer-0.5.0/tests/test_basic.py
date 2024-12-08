# some basic tests

from pathlib import Path
from pptx_renderer import PPTXRenderer
from pptx_renderer.exceptions import RenderError
import pytest
from pptx import Presentation


def mymethod(abc):
    return f"{abc} " * 5


def getimage():
    return "docs/_src/_static/is_it_worth_the_time.png"


mytable = [["a", "b", "c", "d", "e"]] * 10


def test_initiation():
    p = PPTXRenderer("template.pptx")
    assert p.template_path == "template.pptx"


def test_file_not_exist():
    p = PPTXRenderer("nonexistentfile.pptx")
    with pytest.raises(FileNotFoundError):
        p.render("output.pptx")


def test_render():
    p = PPTXRenderer("template.pptx")
    with pytest.raises(RenderError):
        p.render(
            "output.pptx",
            {"mymethod": mymethod, "getimage": getimage, "mytable": mytable},
        )


def test_render_skip_failed():
    p = PPTXRenderer("template.pptx")
    with pytest.warns(UserWarning):
        p.render(
            "output.pptx",
            {"mymethod": mymethod, "getimage": getimage, "mytable": mytable},
            skip_failed=True,
        )
    assert Path("output.pptx").exists()

def test_multi_variable_textbox():
    p = PPTXRenderer("template.pptx")
    p.render(
        "output.pptx",
        {"mymethod": mymethod, "getimage": getimage, "mytable": mytable},
        skip_failed=True,
    )
    prs = Presentation("output.pptx")
    slide = prs.slides[0]
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.startswith("The value of X is"):
            assert shape.text.strip() == "The value of X is 1 and the value of Y is 10."

def test_loop_group():
    template = Presentation("template.pptx")
    slide = template.slides[0]
    title = slide.shapes.title
    title.text = "Loop Group Test - {{{i}}}"
    template.save("template-loop-group.pptx")
    p = PPTXRenderer("template-loop-group.pptx")
    someval = "hello"
    def mymethod(abc):
        return f"{abc} " * 5
    p.render(
        "output.pptx",
        {
            "someval": someval,
            "mymethod": mymethod,
            "myimage": getimage(),
            "mytable": mytable,
        },
        loop_groups=[{"start": 0, "end": 0, "variable": "i", "iterable": [1, 2, 3]}],
    )
    prs = Presentation("output.pptx")
    assert len(prs.slides) == 3
    slide = prs.slides[0]
    assert slide.shapes.title.text == "Loop Group Test - 1"