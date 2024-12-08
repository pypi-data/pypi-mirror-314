from pathlib import Path
from pptx import Presentation
from subprocess import run
from .fixtures import input_json, input_template


def test_basic_cli():
    """Test the basic cli"""
    run(["pptx-renderer", "template.pptx", "output_cli.pptx"])
    assert Path("output_cli.pptx").exists()
    Path("output_cli.pptx").unlink()


def test_with_input(input_template, input_json):
    """Test passing an input"""
    run(
        [
            "pptx-renderer",
            str(input_template),
            "output_cli.pptx",
            "--input",
            str(input_json),
        ]
    )
    assert Path("output_cli.pptx").exists()
    p = Presentation("output_cli.pptx")
    # check if the first slide contains a textbox with the text "hello 1"
    print(p.slides[0].shapes[0].text_frame.text)
    print(p.slides[0].shapes[1].text_frame.text)
    assert "hello 1" in p.slides[0].shapes[0].text_frame.text
    assert "hello 2" in p.slides[0].shapes[1].text_frame.text
    Path("output_cli.pptx").unlink()
