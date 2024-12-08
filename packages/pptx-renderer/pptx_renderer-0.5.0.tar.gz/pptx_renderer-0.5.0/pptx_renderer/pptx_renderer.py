"""Main Module"""

import re
from os import PathLike
from pathlib import Path
from typing import Any, Dict, Optional, Union, Callable, List
from warnings import warn as warning
from functools import partial
from . import plugins

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .exceptions import RenderError
from .utils import fix_quotes, para_text_replace, copy_slide, clear_presentation

PLUGINS = [plugins.image, plugins.video, plugins.table]

class PPTXRenderer:
    """PPTX Renderer class

    This class is used to render a PPTX template by replacing python statements
    with the result of evaluating the python statements.

    Attributes:
        template_path (str): Path to the PPTX template.
    """

    def __init__(self, template_path: Union[str, PathLike]):
        self.template_path = template_path
        self.plugins = {}
        self.namespace = {}
        for plugin in PLUGINS:
            self.register_plugin(plugin.__name__, plugin)

    def register_plugin(self, name: str, func: Callable):
        """Register a plugin function.

        The plugin function should take 2 or more arguments. The first argument
        is the result of evaluating the python statement. The second argument is
        a dictionary containing the following keys:
        - result: The result of evaluating the python statement
        - presentation: The output pptx presentation object
        - shape: The pptx shape object where the placeholder is present
        - slide: The pptx slide object where the placeholder is present
        - slide_no: The slide number where the placeholder is present
        The remaining arguments are the arguments passed to the plugin function

        Args:
            name (str): Name of the plugin.
            func (callable): Function to be registered.

        Returns:
            None
        """
        self.plugins[name] = func

    def render(
        self,
        output_path: Union[str, PathLike],
        methods_and_params: Optional[Dict[str, Any]] = None,
        skip_failed: bool = False,
        loop_groups: Optional[List[Dict[str, Any]]] = None,
    ) -> None:
        """Render PPTXRenderer template and save to output_path.

        Args:
            output_path (str): Path to the output PPTX file.
            methods_and_params (dict, optional): Dictionary of methods and parameters
                to be used in the template. Defaults to None.
            skip_failed (bool, optional): Dont raise an error if some of the
                statements failed to render. Defaults to False.
            loop_groups (list, optional): List of dictionaries containing the
                following keys:
                - start: Slide number where the loop starts
                - end: Slide number where the loop ends
                - variable: Variable name to be used in the loop
                - iterable: Iterable to loop over
                Defaults to None.

        Returns:
            None
        """
        self.template_path = str(self.template_path)
        if not Path(self.template_path).exists():
            raise (FileNotFoundError(f"{self.template_path} not found"))
        template_ppt = Presentation(self.template_path)
        if not methods_and_params:
            methods_and_params = {}
        self.namespace.update(methods_and_params)
        if not loop_groups:
            loop_groups = []

        #function to recurse through list of shapes
        #hand off to function to process text frames and tables
        #recurse into handle_shapes function if group of shapes are found
        def handle_shapes(shapes):
            for shape in list(shapes):
                if shape.has_text_frame:
                    handle_text_frame(shape)
                if shape.has_table:
                    handle_table(shape)
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    handle_shapes(shape.shapes)

        def handle_text_frame(shape):
            matches = re.finditer(r"{{{(.*?)}}}", shape.text)
            if not matches:
                return
            for match_assignment in matches:
                parts = match_assignment.group(1).split(":", 1)
                try:
                    result = eval(fix_quotes(parts[0]), self.namespace)
                except Exception as ex:
                    if skip_failed:
                        warning(
                            f"Evaluation of '{parts[0]}' in slide {slide_no+1} failed"
                        )
                        return
                    raise RenderError(
                        f"Failed to evaluate '{parts[0]}' in slide {slide_no+1}."
                    ) from ex
                if len(parts) > 1:
                    namespace = self.namespace.copy()
                    context = {
                        "result": result,
                        "presentation": template_ppt,
                        "shape": shape,
                        "slide": slide,
                        "slide_no": slide_no,
                    }
                    for plugin_name, plugin in self.plugins.items():
                        func = partial(plugin, context)
                        namespace[plugin_name] = func 
                    try:
                        exec(fix_quotes(parts[1]), namespace)
                    except Exception as ex:
                        if skip_failed:
                            warning(
                                f"Failed to render {parts[0]} in slide {slide_no+1}"
                            )
                            return
                        raise RenderError(
                            f"Failed to render {parts[0]} in slide {slide_no+1}"
                        ) from ex
                else:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text_replace(
                            paragraph, match_assignment.group(0), result
                        )
        
        def handle_table(shape):
            for row in shape.table.rows:
                for cell in row.cells:
                    matches = re.finditer(r"{{{(.*)}}}", cell.text)
                    if not matches:
                        continue
                    for match_assignment in matches:
                        parts = match_assignment.group(1).split(":", 1)
                        try:
                            result = eval(fix_quotes(parts[0]), self.namespace)
                        except Exception as ex:
                            if skip_failed:
                                warning(
                                    f"Evaluation of '{parts[0]}' in slide {slide_no+1} failed"
                                )
                                continue
                            raise RenderError(
                                f"Failed to evaluate '{parts[0]}'."
                            ) from ex
                        for paragraph in cell.text_frame.paragraphs:
                            para_text_replace(
                                paragraph, match_assignment.group(0), result
                            )

        output_ppt = Presentation(self.template_path)
        clear_presentation(output_ppt)
        extra_namespace = {}
        slides_managed = []
        for slide_no, slide in enumerate(template_ppt.slides):
            if slide_no in slides_managed:
                # this slide has already been looped over
                continue
            slide_used = False
            for loop_group in loop_groups:
                if slide_no == loop_group["start"]:
                    slide_used = True
                    for variable_value in loop_group["iterable"]:
                        for loop_slide_no in range(
                            loop_group["start"], loop_group["end"] + 1
                        ):
                            if loop_slide_no not in slides_managed:
                                slides_managed.append(loop_slide_no)
                            # get the slide from the template
                            current_slide = template_ppt.slides[loop_slide_no]
                            # add a copy of this slide to output_slides
                            new_slide = copy_slide(template_ppt, output_ppt, current_slide)
                            extra_namespace[output_ppt.slides.index(new_slide)] = {
                                loop_group["variable"]: variable_value
                            }
            if not slide_used:
                # this slide is not part of a loop group
                new_slide = copy_slide(template_ppt, output_ppt, slide)
        for slide_no, slide in enumerate(output_ppt.slides):
            self.namespace.update(extra_namespace.get(slide_no, {}))
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                python_code = re.search(
                    r"```python([\s\S]*)```",
                    fix_quotes(slide.notes_slide.notes_text_frame.text),
                    re.MULTILINE,
                )
                if python_code:
                    exec(python_code.group(1), self.namespace)
            #send list of shapes to handle_shapes function
            handle_shapes(list(slide.shapes))
        output_ppt.save(output_path)
