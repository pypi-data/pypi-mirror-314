import copy
import io

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.shapes.group import GroupShape


def para_text_replace(para, find_string, replace_string):
    """Function to replace text in a paragraph

    This function replaces text in a paragraph while respecting the formatting.

    Args:
        para (pptx.shapes.paragraph.Paragraph): Paragraph to replace text in.
        find_string (str): String to find in the paragraph.
        replace_string (str): String to replace the find_string with.

    Returns:
        None
    """
    find_string = str(find_string)
    replace_string = str(replace_string)
    starting_pos = para.text.find(find_string)
    if starting_pos == -1:
        return  # text not in paragraph
    txt_prev = ""
    for run in para.runs:
        if len(txt_prev) <= starting_pos < len(txt_prev) + len(run.text):
            if run.text.find(find_string) != -1:  # text in run, replace
                run.text = run.text.replace(find_string, replace_string)
                return
            else:  # text no in "run"
                txt_prev = txt_prev + run.text
                run.text = run.text[: starting_pos - len(txt_prev)] + replace_string
        elif starting_pos < len(txt_prev) and starting_pos + len(find_string) >= len(
            txt_prev
        ) + len(run.text):
            txt_prev = txt_prev + run.text
            run.text = ""
        elif (
            len(txt_prev)
            < starting_pos + len(find_string)
            < len(txt_prev) + len(run.text)
        ):
            txt_prev = txt_prev + run.text
            run.text = run.text[starting_pos + len(find_string) - len(txt_prev) :]
        else:
            txt_prev += run.text


def fix_quotes(input_string: str) -> str:
    """Replace unicode quotes (inserted by powerpoint) with ascii quotes.

    Args:
        input_string (str): String to fix quotes in.
    
    Returns:
        str: String with fixed quotes.
    """
    return (
        input_string.replace("’", "'")
        .replace("‘", "'")
        .replace("“", '"')
        .replace("”", '"')
    )

def find_layout_index(prs: Presentation, slide: Slide):
    """Finds the index of the layout of a slide in a presentation

    Args:
        prs (Presentation): Presentation to find the layout index in.
        slide (Slide): Slide to find the layout index of.

    Returns:
        int: Index of the layout of the slide in the presentation.
    """
    layouts_per_master = [list(master.slide_layouts) for master in prs.slide_masters]
    slide_layout = slide.slide_layout
    for outer_index, inner_list in enumerate(layouts_per_master):
        if slide_layout in inner_list:
            inner_index = inner_list.index(slide_layout)
            return outer_index, inner_index
    raise ValueError("Slide layout not found in presentation")

def copy_table(source_shape, target_slide):
    source_table = source_shape.table
    # Get the dimensions of the source table
    rows = len(source_table.rows)
    cols = len(source_table.columns)

    # Add a new table to the target slide with the same dimensions
    left = source_shape.left
    top = source_shape.top
    width = source_shape.width
    height = source_shape.height
    target_table = target_slide.shapes.add_table(
        rows, cols, left, top, width, height
    ).table

    # Copy the content and formatting of each cell
    for row_idx in range(rows):
        for col_idx in range(cols):
            source_cell = source_table.cell(row_idx, col_idx)
            target_cell = target_table.cell(row_idx, col_idx)

            copy_text_frame(source_cell.text_frame, target_cell.text_frame)


def copy_shapes(shapes, target_slide, ignore_placeholders=True):
    """Copy shapes from one slide to another

    Args:
        shapes (list): List of shapes to copy.
        target_slide (Slide): Slide to copy the shapes to.
        ignore_placeholders (bool, optional): Ignore placeholders when copying shapes. Defaults to True.

    Returns:
        None
    """
    # Copy all existing shapes
    for shape in shapes:
        if shape.is_placeholder and ignore_placeholders:
            continue
        if isinstance(shape, GroupShape):
            group = target_slide.shapes.add_group_shape()
            group.name = shape.name
            group.left = shape.left
            group.top = shape.top
            group.width = shape.width
            group.height = shape.height
            group.rotation = shape.rotation

            # Recursive copy of contents
            copy_shapes(shape.shapes, group)

            # Fix offset
            cur_el = group._element.xpath(".//p:grpSpPr")[0]
            ref_el = shape._element.xpath(".//p:grpSpPr")[0]
            parent = cur_el.getparent()
            parent.insert(parent.index(cur_el) + 1, copy.deepcopy(ref_el))
            parent.remove(cur_el)

            result = group
        elif hasattr(shape, "image"):
            # Get image contents
            content = io.BytesIO(shape.image.blob)
            result = target_slide.shapes.add_picture(
                content, shape.left, shape.top, shape.width, shape.height
            )
            result.name = shape.name
            result.crop_left = shape.crop_left
            result.crop_right = shape.crop_right
            result.crop_top = shape.crop_top
            result.crop_bottom = shape.crop_bottom
        elif hasattr(shape, "has_chart") and shape.has_chart:
            pass
        elif hasattr(shape, "has_table") and shape.has_table:
            # copy table style
            copy_table(shape, target_slide)
        else:
            newel = copy.deepcopy(shape.element)
            target_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
            result = target_slide.shapes[-1]
            if shape.has_text_frame:
                copy_text_frame(shape.text_frame, result.text_frame)


def copy_text_frame(source_text_frame, target_text_frame):
    try:
        target_text_frame.clear()  # Clear any existing text
    except AttributeError:
        pass  # Text frame is not clearable
    if not hasattr(source_text_frame, "paragraphs"):
        return
    for i, paragraph in enumerate(source_text_frame.paragraphs):
        if i == 0:
            # First paragraph is already there
            new_paragraph = target_text_frame.paragraphs[0]
        else:
            new_paragraph = target_text_frame.add_paragraph()
        new_paragraph.level = paragraph.level

        # Copy paragraph formatting
        new_paragraph.font.bold = paragraph.font.bold
        new_paragraph.font.italic = paragraph.font.italic
        new_paragraph.font.underline = paragraph.font.underline
        new_paragraph.font.size = paragraph.font.size
        if paragraph.font.color and hasattr(paragraph.font.color, "rgb"):
            new_paragraph.font.color.rgb = paragraph.font.color.rgb
        new_paragraph.alignment = paragraph.alignment

        for run in paragraph.runs:
            new_run = new_paragraph.add_run()
            new_run.text = run.text

            # Copy run formatting
            new_run.font.bold = run.font.bold
            new_run.font.italic = run.font.italic
            new_run.font.underline = run.font.underline
            new_run.font.size = run.font.size
            if run.font.color and hasattr(run.font.color, "rgb") and run.font.color.rgb:
                new_run.font.color.rgb = run.font.color.rgb
            elif run.font.color and hasattr(run.font.color, "theme_color"):
                new_run.font.color.theme_color = run.font.color.theme_color
            if run.font.color and hasattr(run.font.color, "brightness") and run.font.color.brightness:
                new_run.font.color.brightness = run.font.color.brightness
            if (
                hasattr(run, "hyperlink")
                and run.hyperlink
                and hasattr(run.hyperlink, "address")
            ):
                new_run.hyperlink.address = run.hyperlink.address

def copy_slide(source_ppt: Presentation, target_ppt: Presentation, slide: Slide) -> Slide:
    """Duplicate each slide in prs2 and "moves" it into prs1.
    Adds slides to the end of the presentation

    Args:
        source_ppt (Presentation): Source presentation.
        target_ppt (Presentation): Target presentation.
        slide (Slide): Slide to copy.

    Returns:
        Slide: Slide that was copied
    """
    master_index, layout_index = find_layout_index(source_ppt, slide)
    target_ppt_layouts = [list(master.slide_layouts) for master in target_ppt.slide_masters]
    new_slide = target_ppt.slides.add_slide(target_ppt_layouts[master_index][layout_index])
    for new_ph, old_ph in zip(new_slide.placeholders, slide.placeholders):
        new_ph.height = old_ph.height
        new_ph.width = old_ph.width
        new_ph.left = old_ph.left
        new_ph.top = old_ph.top
        if not old_ph.has_text_frame or not new_ph.has_text_frame:
            continue
        if not hasattr(old_ph.text_frame, "paragraphs") or not hasattr(
            new_ph.text_frame, "paragraphs"
        ):
            continue
        copy_text_frame(old_ph.text_frame, new_ph.text_frame)
    copy_shapes(slide.shapes, new_slide, ignore_placeholders=True)
    # copy slide notes
    if slide.notes_slide:
        for new_ph, old_ph in zip(
            new_slide.notes_slide.placeholders, slide.notes_slide.placeholders
        ):
            copy_text_frame(old_ph.text_frame, new_ph.text_frame)
        new_slide_notes = new_slide.notes_slide
        for shape in slide.notes_slide.placeholders:
            if shape.is_placeholder:
                continue
            newel = copy.deepcopy(shape.element)
            new_slide_notes.shapes._spTree.insert_element_before(newel, "p:extLst")
    return new_slide

def clear_presentation(prs: Presentation):
    """Clears all slides from a presentation"""
    for i in range(len(prs.slides)-1, -1, -1): 
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]